"""
Chargement des documents source et détection des fichiers déjà indexés dans Chroma.

Formats pris en charge : PDF (.pdf), Word (.docx), PowerPoint (.ppt, .pptx).
Chaque document chargé est enrichi de métadonnées (nom de fichier, type, numéro de page)
afin de pouvoir citer les sources dans les réponses de la chaîne RAG.
"""
import logging
from pathlib import Path

from langchain_community.document_loaders import Docx2txtLoader, PyPDFLoader
from langchain_core.documents import Document
from pptx import Presentation

logger = logging.getLogger(__name__)


def _load_pptx(filepath: Path) -> list:
    """Extrait le texte d'un PPTX slide par slide (sans dépendance à unstructured/spaCy).

    Les slides vides sont ignorées.
    """
    prs = Presentation(str(filepath))
    docs = []
    for slide in prs.slides:
        texte = "\n".join(
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
        if texte:
            docs.append(Document(page_content=texte))
    return docs


LOADERS = {
    ".pdf": PyPDFLoader,
    ".docx": Docx2txtLoader,
    ".ppt": _load_pptx,
    ".pptx": _load_pptx,
}


def get_indexed_filenames(vectordb) -> set:
    """Retourne les noms de fichiers déjà présents dans le vector store."""
    resultats = vectordb.get(include=["metadatas"])
    return {m["fichier"] for m in resultats["metadatas"] if "fichier" in m}


def get_indexed_files_info(vectordb) -> dict:
    """Retourne un dict {nom_fichier: nombre_de_chunks} pour tous les fichiers indexés."""
    resultats = vectordb.get(include=["metadatas"])
    counts = {}
    for m in resultats["metadatas"]:
        nom = m.get("fichier")
        if nom:
            counts[nom] = counts.get(nom, 0) + 1
    return counts


def delete_document(vectordb, filename: str, doc_dir: Path) -> None:
    """Supprime un document du vector store et du disque.

    Retire tous les chunks associés au fichier dans Chroma, puis supprime
    le fichier physique dans doc_dir s'il existe.

    Args:
        vectordb: Instance Chroma connectée au vector store persisté.
        filename: Nom du fichier à supprimer (ex: "rapport.pdf").
        doc_dir: Répertoire où le fichier est stocké sur disque.
    """
    resultats = vectordb.get(where={"fichier": filename})
    ids = resultats.get("ids", [])
    if ids:
        vectordb.delete(ids=ids)
        logger.info("Supprimé de Chroma : %s (%d chunks)", filename, len(ids))

    fichier_path = doc_dir / filename
    if fichier_path.exists():
        fichier_path.unlink()
        logger.info("Fichier supprimé du disque : %s", fichier_path)


def load_file(fichier: Path) -> list:
    """Charge un fichier et enrichit chaque page/slide avec les métadonnées de source.

    Sélectionne le loader selon l'extension (doit être dans LOADERS), puis
    ajoute `fichier`, `type` et `page_number` à chaque Document.
    """
    loader = LOADERS[fichier.suffix.lower()]
    # _load_pptx est une fonction directe ; les loaders LangChain sont des classes
    if callable(loader) and not isinstance(loader, type):
        docs = loader(fichier)
    else:
        docs = loader(str(fichier)).load()

    for i, doc in enumerate(docs, start=1):
        doc.metadata["fichier"] = fichier.name
        doc.metadata["type"] = fichier.suffix.lower().lstrip(".")
        doc.metadata["page_number"] = i

    return docs


def load_new_documents(repertoire: Path, deja_indexes: set) -> list:
    """Charge les fichiers de `repertoire` absents de `deja_indexes`.

    Parcourt récursivement le répertoire. Les erreurs sur un fichier sont
    loggées sans interrompre le traitement des suivants.
    """
    fichiers_trouves = [f for f in repertoire.rglob("*") if f.suffix.lower() in LOADERS]

    if not fichiers_trouves:
        logger.warning("Aucun fichier PDF/DOCX/PPTX trouvé dans %s", repertoire)
        return []

    nouveaux = [f for f in fichiers_trouves if f.name not in deja_indexes]
    deja_vus = [f for f in fichiers_trouves if f.name in deja_indexes]

    logger.info("%d fichier(s) trouvé(s) dans %s", len(fichiers_trouves), repertoire)
    logger.info("Déjà indexés : %d | Nouveaux : %d", len(deja_vus), len(nouveaux))

    for fichier in deja_vus:
        logger.debug("Ignoré (déjà indexé) : %s", fichier.name)

    tous_les_docs = []
    for fichier in nouveaux:
        try:
            docs = load_file(fichier)
            tous_les_docs.extend(docs)
            logger.info("%s -> %d page(s)/slide(s) chargé(s)", fichier.name, len(docs))
        except Exception:
            logger.exception("Erreur lors du chargement de %s", fichier.name)

    return tous_les_docs
